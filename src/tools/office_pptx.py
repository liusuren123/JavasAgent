"""PowerPoint 操作子模块。

提供 PowerPoint (.pptx) 文件的创建、读取、编辑能力。
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

from loguru import logger

from src.utils.path_safety import PathSafetyError, safe_resolve_path


# ======================================================================
# 公共辅助
# ======================================================================


def _safe_path(workspace: Path, user_path: str, *, allow_create_parents: bool = False) -> Path:
    """安全解析用户路径。"""
    return safe_resolve_path(workspace, user_path, allow_create_parents=allow_create_parents)


def _check_pptx_import() -> str | None:
    """检查 python-pptx 是否已安装。"""
    try:
        import pptx  # noqa: F401
        return None
    except ImportError:
        return "python-pptx 未安装，请运行: pip install python-pptx"


# ======================================================================
# 操作实现
# ======================================================================


async def read_pptx(workspace: Path, params: dict) -> dict:
    """读取 PowerPoint 演示文稿内容。

    Params:
        path: 文件路径
        max_slides: 最大幻灯片数（默认 50）
    """
    try:
        path = _safe_path(workspace, params["path"])
    except PathSafetyError as e:
        return {"error": str(e)}

    if not path.exists():
        return {"error": f"文件不存在: {path}"}

    err = _check_pptx_import()
    if err:
        return {"error": err}

    try:
        from pptx import Presentation

        prs = Presentation(str(path))
        max_slides = params.get("max_slides", 50)

        slides: list[dict[str, Any]] = []
        for i, slide in enumerate(prs.slides):
            if i >= max_slides:
                break
            shapes: list[dict[str, str]] = []
            try:
                for shape in slide.shapes:
                    try:
                        if shape.has_text_frame:
                            text = shape.text_frame.text.strip()
                            if text:
                                shapes.append({
                                    "name": shape.name,
                                    "type": "text",
                                    "text": text[:500],
                                })
                    except Exception:
                        continue
            except Exception as shape_err:
                logger.warning(f"幻灯片 {i + 1} shape 遍历异常: {shape_err}")
            slides.append({"slide_number": i + 1, "shapes": shapes})

        logger.info(f"读取 PPT: {path} ({len(slides)} 张幻灯片)")
        return {
            "path": str(path),
            "slides": slides,
            "slide_count": len(prs.slides),
            "width": prs.slide_width,
            "height": prs.slide_height,
        }
    except Exception as e:
        logger.error(f"读取 PPT 失败: {e}")
        return {"error": f"读取失败: {e}"}


async def create_pptx(workspace: Path, params: dict) -> dict:
    """创建 PowerPoint 演示文稿。

    Params:
        path: 文件路径
        title: 标题（可选）
        slides: 幻灯片列表 [{"title": "标题", "body": "内容"}]
    """
    try:
        path = _safe_path(workspace, params["path"], allow_create_parents=True)
    except PathSafetyError as e:
        return {"error": str(e)}

    err = _check_pptx_import()
    if err:
        return {"error": err}

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation()

        slide_data = params.get("slides", [])
        for slide_info in slide_data:
            blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
            slide = prs.slides.add_slide(blank_layout)

            # 手动添加标题文本框
            title_text = slide_info.get("title", "")
            if title_text:
                title_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(0.3), Inches(9), Inches(1)
                )
                tf = title_box.text_frame
                tf.text = title_text
                for para in tf.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(28)
                        run.font.bold = True

            # 手动添加内容文本框
            body_text = slide_info.get("body", "")
            if body_text:
                body_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(1.5), Inches(9), Inches(5)
                )
                body_box.text_frame.text = body_text

        path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(path))

        logger.info(f"创建 PPT: {path} ({len(slide_data)} 张幻灯片)")
        return {
            "path": str(path),
            "created": True,
            "slide_count": len(slide_data),
        }
    except Exception as e:
        logger.error(f"创建 PPT 失败: {e}")
        return {"error": f"创建失败: {e}"}


async def edit_pptx(workspace: Path, params: dict) -> dict:
    """编辑已有 PowerPoint 演示文稿。

    支持的操作（通过 ``operations`` 列表指定，按顺序执行）：

    1. **add_slide** — 添加新幻灯片
       ``{"action": "add_slide", "title": "标题", "body": "正文"}``

    2. **delete_slide** — 删除幻灯片
       ``{"action": "delete_slide", "slide_index": 0}``  (0-based)

    3. **replace_text** — 替换文本（全局搜索替换）
       ``{"action": "replace_text", "search": "旧文本", "replace": "新文本"}``

    4. **add_textbox** — 在指定幻灯片上添加文本框
       ``{"action": "add_textbox", "slide_index": 0, "text": "内容",
          "left": 1.0, "top": 2.0, "width": 5.0, "height": 1.0,
          "font_size": 18, "bold": true}``

    5. **add_image** — 在指定幻灯片上添加图片
       ``{"action": "add_image", "slide_index": 0, "image_path": "img.png",
          "left": 1.0, "top": 2.0, "width": 4.0, "height": 3.0}``

    Params:
        path: 文件路径
        operations: 操作列表（见上文）
    """
    try:
        path = _safe_path(workspace, params["path"])
    except PathSafetyError as e:
        return {"error": str(e)}

    if not path.exists():
        return {"error": f"文件不存在: {path}"}

    err = _check_pptx_import()
    if err:
        return {"error": err}

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation(str(path))
        operations = params.get("operations", [])
        results: list[dict[str, Any]] = []

        for op in operations:
            action = op.get("action")
            try:
                if action == "add_slide":
                    res = _op_add_slide(prs, op)
                    results.append(res)

                elif action == "delete_slide":
                    res = _op_delete_slide(prs, op)
                    results.append(res)

                elif action == "replace_text":
                    res = _op_replace_text(prs, op)
                    results.append(res)

                elif action == "add_textbox":
                    res = _op_add_textbox(prs, workspace, op)
                    results.append(res)

                elif action == "add_image":
                    res = _op_add_image(prs, workspace, op)
                    results.append(res)

                else:
                    results.append({"action": action, "error": f"未知操作: {action}"})
            except Exception as e:
                results.append({"action": action, "error": str(e)})

        path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(path))

        logger.info(f"编辑 PPT: {path} ({len(operations)} 个操作)")
        return {
            "path": str(path),
            "operations": len(operations),
            "details": results,
        }
    except Exception as e:
        logger.error(f"编辑 PPT 失败: {e}")
        return {"error": f"编辑失败: {e}"}


# ======================================================================
# 编辑操作实现
# ======================================================================


def _op_add_slide(prs: Any, op: dict) -> dict[str, Any]:
    """添加新幻灯片。"""
    from pptx.util import Inches, Pt

    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(blank_layout)

    title_text = op.get("title", "")
    if title_text:
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
        tf = title_box.text_frame
        tf.text = title_text
        for para in tf.paragraphs:
            for run in para.runs:
                run.font.size = Pt(28)
                run.font.bold = True

    body_text = op.get("body", "")
    if body_text:
        body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
        body_box.text_frame.text = body_text

    return {"action": "add_slide", "status": "ok"}


def _op_delete_slide(prs: Any, op: dict) -> dict[str, Any]:
    """删除幻灯片（按 0-based 索引）。"""
    idx = op.get("slide_index", -1)
    if idx < 0 or idx >= len(prs.slides):
        return {"action": "delete_slide", "error": f"无效索引: {idx}，共 {len(prs.slides)} 张"}

    # 使用 python-pptx 推荐的方式：直接删除对应的 XML 节点
    slide = prs.slides[idx]
    rId = prs.slides._sldIdLst[idx].attrib[
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    ]
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[idx])

    return {"action": "delete_slide", "status": "ok", "deleted_index": idx}


def _op_replace_text(prs: Any, op: dict) -> dict[str, Any]:
    """全局搜索替换文本。"""
    search = op.get("search", "")
    replace = op.get("replace", "")
    if not search:
        return {"action": "replace_text", "error": "search 不能为空"}

    count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if search in run.text:
                            run.text = run.text.replace(search, replace)
                            count += 1

    return {"action": "replace_text", "status": "ok", "replacements": count}


def _op_add_textbox(prs: Any, workspace: Path, op: dict) -> dict[str, Any]:
    """在指定幻灯片上添加文本框。"""
    from pptx.util import Inches, Pt

    idx = op.get("slide_index", 0)
    if idx < 0 or idx >= len(prs.slides):
        return {"action": "add_textbox", "error": f"无效幻灯片索引: {idx}"}

    slide = prs.slides[idx]
    left = Inches(op.get("left", 0.5))
    top = Inches(op.get("top", 1.5))
    width = Inches(op.get("width", 9.0))
    height = Inches(op.get("height", 1.0))

    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.text = op.get("text", "")

    font_size = op.get("font_size")
    bold = op.get("bold", False)
    if font_size or bold:
        for para in tf.paragraphs:
            for run in para.runs:
                if font_size:
                    run.font.size = Pt(font_size)
                if bold:
                    run.font.bold = True

    return {"action": "add_textbox", "status": "ok"}


def _op_add_image(prs: Any, workspace: Path, op: dict) -> dict[str, Any]:
    """在指定幻灯片上添加图片。"""
    from pptx.util import Inches

    idx = op.get("slide_index", 0)
    if idx < 0 or idx >= len(prs.slides):
        return {"action": "add_image", "error": f"无效幻灯片索引: {idx}"}

    # 解析图片路径（相对于 workspace）
    img_path_str = op.get("image_path", "")
    img_path = safe_resolve_path(workspace, img_path_str)
    if not img_path.exists():
        return {"action": "add_image", "error": f"图片不存在: {img_path}"}

    slide = prs.slides[idx]
    left = Inches(op.get("left", 1.0))
    top = Inches(op.get("top", 2.0))
    width = Inches(op.get("width", 4.0))
    height = Inches(op.get("height", 3.0))

    slide.shapes.add_picture(str(img_path), left, top, width, height)

    return {"action": "add_image", "status": "ok"}
