"""办公自动化工具集。

提供 Word/Excel/PPT 文档操作、邮件处理等办公自动化能力。
基于 python-docx / openpyxl / python-pptx 等库实现。
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

from loguru import logger

from src.utils.command import run_command
from src.utils.path_safety import PathSafetyError, safe_resolve_path


class OfficeOps:
    """办公自动化工具集。

    支持 Word (.docx)、Excel (.xlsx)、PowerPoint (.pptx) 文档的
    创建、读取、编辑操作。

    Usage::

        office = OfficeOps(workspace="/path/to/workspace")
        # 读取 Word 文档
        result = await office.execute("read_docx", {"path": "report.docx"})
        # 创建 Excel 表格
        result = await office.execute("create_xlsx", {
            "path": "data.xlsx",
            "headers": ["姓名", "年龄"],
            "rows": [["张三", 25]],
        })
    """

    def __init__(self, workspace: str | None = None) -> None:
        self._workspace = Path(workspace) if workspace else Path.cwd()

    async def execute(self, action: str, params: dict[str, Any]) -> Any:
        """执行办公自动化操作。

        Args:
            action: 操作类型
            params: 操作参数

        Returns:
            操作结果字典
        """
        handlers = {
            # Word 操作
            "read_docx": self._read_docx,
            "create_docx": self._create_docx,
            "append_docx": self._append_docx,
            # Excel 操作
            "read_xlsx": self._read_xlsx,
            "create_xlsx": self._create_xlsx,
            "append_xlsx": self._append_xlsx,
            # PowerPoint 操作
            "read_pptx": self._read_pptx,
            "create_pptx": self._create_pptx,
            # PDF
            "read_pdf_text": self._read_pdf_text,
        }

        handler = handlers.get(action)
        if handler is None:
            logger.error(f"未知操作: {action}")
            return {
                "error": f"未知操作: {action}",
                "available_actions": sorted(handlers.keys()),
            }

        return await handler(params)

    def _safe_path(self, user_path: str, *, allow_create_parents: bool = False) -> Path:
        """安全解析用户路径，防止路径遍历。"""
        return safe_resolve_path(
            self._workspace,
            user_path,
            allow_create_parents=allow_create_parents,
        )

    # ------------------------------------------------------------------
    # Word (.docx) 操作
    # ------------------------------------------------------------------

    async def _read_docx(self, params: dict) -> dict:
        """读取 Word 文档内容。

        Params:
            path: 文件路径（相对于 workspace）
            max_paragraphs: 最大读取段落数（默认 500）
        """
        try:
            path = self._safe_path(params["path"])
        except PathSafetyError as e:
            return {"error": str(e)}

        if not path.exists():
            return {"error": f"文件不存在: {path}"}

        try:
            from docx import Document
        except ImportError:
            return {"error": "python-docx 未安装，请运行: pip install python-docx"}

        try:
            doc = Document(str(path))
            max_para = params.get("max_paragraphs", 500)

            paragraphs: list[dict[str, Any]] = []
            for i, para in enumerate(doc.paragraphs[:max_para]):
                if para.text.strip():
                    paragraphs.append({
                        "index": i,
                        "style": para.style.name if para.style else "Normal",
                        "text": para.text,
                    })

            # 提取表格内容
            tables: list[list[list[str]]] = []
            for table in doc.tables:
                rows_data: list[list[str]] = []
                for row in table.rows:
                    rows_data.append([cell.text for cell in row.cells])
                tables.append(rows_data)

            logger.info(f"读取 Word 文档: {path} ({len(paragraphs)} 段落, {len(tables)} 表格)")
            return {
                "path": str(path),
                "paragraphs": paragraphs,
                "paragraph_count": len(paragraphs),
                "tables": tables,
                "table_count": len(tables),
            }
        except Exception as e:
            logger.error(f"读取 Word 文档失败: {e}")
            return {"error": f"读取失败: {e}"}

    async def _create_docx(self, params: dict) -> dict:
        """创建 Word 文档。

        Params:
            path: 文件路径
            title: 文档标题（可选）
            paragraphs: 段落文本列表
            headings: 标题列表 [{"level": 1, "text": "标题"}]（可选）
        """
        try:
            path = self._safe_path(params["path"], allow_create_parents=True)
        except PathSafetyError as e:
            return {"error": str(e)}

        try:
            from docx import Document
            from docx.shared import Pt
        except ImportError:
            return {"error": "python-docx 未安装，请运行: pip install python-docx"}

        try:
            doc = Document()

            # 添加标题
            title = params.get("title")
            if title:
                doc.add_heading(title, level=0)

            # 添加标题列表
            for heading in params.get("headings", []):
                level = min(heading.get("level", 1), 4)
                doc.add_heading(heading["text"], level=level)

            # 添加段落
            for text in params.get("paragraphs", []):
                doc.add_paragraph(text)

            path.parent.mkdir(parents=True, exist_ok=True)
            doc.save(str(path))

            logger.info(f"创建 Word 文档: {path}")
            return {"path": str(path), "created": True}
        except Exception as e:
            logger.error(f"创建 Word 文档失败: {e}")
            return {"error": f"创建失败: {e}"}

    async def _append_docx(self, params: dict) -> dict:
        """追加内容到现有 Word 文档。

        Params:
            path: 文件路径
            paragraphs: 要追加的段落文本列表
            headings: 要追加的标题列表（可选）
        """
        try:
            path = self._safe_path(params["path"])
        except PathSafetyError as e:
            return {"error": str(e)}

        if not path.exists():
            return {"error": f"文件不存在: {path}"}

        try:
            from docx import Document
        except ImportError:
            return {"error": "python-docx 未安装，请运行: pip install python-docx"}

        try:
            doc = Document(str(path))

            for heading in params.get("headings", []):
                level = min(heading.get("level", 1), 4)
                doc.add_heading(heading["text"], level=level)

            for text in params.get("paragraphs", []):
                doc.add_paragraph(text)

            doc.save(str(path))

            added_count = len(params.get("paragraphs", [])) + len(params.get("headings", []))
            logger.info(f"追加 Word 文档: {path} (+{added_count} 段)")
            return {"path": str(path), "appended": added_count}
        except Exception as e:
            logger.error(f"追加 Word 文档失败: {e}")
            return {"error": f"追加失败: {e}"}

    # ------------------------------------------------------------------
    # Excel (.xlsx) 操作
    # ------------------------------------------------------------------

    async def _read_xlsx(self, params: dict) -> dict:
        """读取 Excel 文件内容。

        Params:
            path: 文件路径
            sheet: 工作表名称（默认第一个）
            max_rows: 最大行数（默认 1000）
            include_formulas: 是否包含公式（默认 False）
        """
        try:
            path = self._safe_path(params["path"])
        except PathSafetyError as e:
            return {"error": str(e)}

        if not path.exists():
            return {"error": f"文件不存在: {path}"}

        try:
            import openpyxl
        except ImportError:
            return {"error": "openpyxl 未安装，请运行: pip install openpyxl"}

        try:
            wb = openpyxl.load_workbook(str(path), data_only=not params.get("include_formulas", False))
            sheet_name = params.get("sheet")
            ws = wb[sheet_name] if sheet_name else wb.active

            max_rows = params.get("max_rows", 1000)
            rows: list[list[Any]] = []
            for row in ws.iter_rows(max_row=min(ws.max_row, max_rows), values_only=True):
                rows.append(list(row))

            # 获取所有 sheet 名称
            sheet_names = wb.sheetnames
            wb.close()

            logger.info(f"读取 Excel: {path} ({ws.title}, {len(rows)} 行)")
            return {
                "path": str(path),
                "sheet": ws.title,
                "sheet_names": sheet_names,
                "rows": rows,
                "row_count": len(rows),
                "column_count": ws.max_column,
            }
        except KeyError:
            return {"error": f"工作表不存在: {sheet_name}，可用: {wb.sheetnames}"}
        except Exception as e:
            logger.error(f"读取 Excel 失败: {e}")
            return {"error": f"读取失败: {e}"}

    async def _create_xlsx(self, params: dict) -> dict:
        """创建 Excel 文件。

        Params:
            path: 文件路径
            sheet: 工作表名称（默认 "Sheet1"）
            headers: 列标题列表
            rows: 数据行列表（二维数组）
            column_widths: 列宽列表（可选）
        """
        try:
            path = self._safe_path(params["path"], allow_create_parents=True)
        except PathSafetyError as e:
            return {"error": str(e)}

        try:
            import openpyxl
        except ImportError:
            return {"error": "openpyxl 未安装，请运行: pip install openpyxl"}

        try:
            wb = openpyxl.Workbook()
            sheet_name = params.get("sheet", "Sheet1")
            ws = wb.active
            ws.title = sheet_name

            headers = params.get("headers", [])
            rows = params.get("rows", [])

            # 写入表头
            if headers:
                ws.append(headers)
                # 表头加粗
                from openpyxl.styles import Font

                for cell in ws[1]:
                    cell.font = Font(bold=True)

            # 写入数据行
            for row in rows:
                ws.append(row)

            # 设置列宽
            widths = params.get("column_widths")
            if widths:
                for i, w in enumerate(widths, 1):
                    ws.column_dimensions[openpyxl.utils.get_column_letter(i)].width = w
            elif headers:
                # 自动调整列宽（基于表头长度估算）
                for i, h in enumerate(headers, 1):
                    ws.column_dimensions[openpyxl.utils.get_column_letter(i)].width = max(
                        len(str(h)) * 2 + 4, 12
                    )

            path.parent.mkdir(parents=True, exist_ok=True)
            wb.save(str(path))
            wb.close()

            total_rows = len(rows) + (1 if headers else 0)
            logger.info(f"创建 Excel: {path} ({total_rows} 行, {len(headers)} 列)")
            return {
                "path": str(path),
                "created": True,
                "sheet": sheet_name,
                "rows": total_rows,
                "columns": len(headers),
            }
        except Exception as e:
            logger.error(f"创建 Excel 失败: {e}")
            return {"error": f"创建失败: {e}"}

    async def _append_xlsx(self, params: dict) -> dict:
        """追加数据到现有 Excel 文件。

        Params:
            path: 文件路径
            rows: 要追加的数据行列表
            sheet: 工作表名称（默认活动工作表）
        """
        try:
            path = self._safe_path(params["path"])
        except PathSafetyError as e:
            return {"error": str(e)}

        if not path.exists():
            return {"error": f"文件不存在: {path}"}

        try:
            import openpyxl
        except ImportError:
            return {"error": "openpyxl 未安装，请运行: pip install openpyxl"}

        try:
            wb = openpyxl.load_workbook(str(path))
            sheet_name = params.get("sheet")
            ws = wb[sheet_name] if sheet_name else wb.active

            rows = params.get("rows", [])
            for row in rows:
                ws.append(row)

            wb.save(str(path))
            wb.close()

            logger.info(f"追加 Excel: {path} (+{len(rows)} 行)")
            return {"path": str(path), "appended_rows": len(rows)}
        except Exception as e:
            logger.error(f"追加 Excel 失败: {e}")
            return {"error": f"追加失败: {e}"}

    # ------------------------------------------------------------------
    # PowerPoint (.pptx) 操作
    # ------------------------------------------------------------------

    async def _read_pptx(self, params: dict) -> dict:
        """读取 PowerPoint 演示文稿内容。

        Params:
            path: 文件路径
            max_slides: 最大幻灯片数（默认 50）
        """
        try:
            path = self._safe_path(params["path"])
        except PathSafetyError as e:
            return {"error": str(e)}

        if not path.exists():
            return {"error": f"文件不存在: {path}"}

        try:
            from pptx import Presentation
        except ImportError:
            return {"error": "python-pptx 未安装，请运行: pip install python-pptx"}

        try:
            prs = Presentation(str(path))
            max_slides = params.get("max_slides", 50)

            slides: list[dict[str, Any]] = []
            for i, slide in enumerate(prs.slides):
                if i >= max_slides:
                    break
                shapes: list[dict[str, str]] = []
                try:
                    for shape in slide.shapes:
                        try:
                            if shape.has_text_frame:
                                text = shape.text_frame.text.strip()
                                if text:
                                    shapes.append({
                                        "name": shape.name,
                                        "type": "text",
                                        "text": text[:500],
                                    })
                        except Exception:
                            # 跳过无法读取的 shape
                            continue
                except Exception as shape_err:
                    logger.warning(f"幻灯片 {i + 1} shape 遍历异常: {shape_err}")
                slides.append({"slide_number": i + 1, "shapes": shapes})

            logger.info(f"读取 PPT: {path} ({len(slides)} 张幻灯片)")
            return {
                "path": str(path),
                "slides": slides,
                "slide_count": len(prs.slides),
                "width": prs.slide_width,
                "height": prs.slide_height,
            }
        except Exception as e:
            logger.error(f"读取 PPT 失败: {e}")
            return {"error": f"读取失败: {e}"}

    async def _create_pptx(self, params: dict) -> dict:
        """创建 PowerPoint 演示文稿。

        Params:
            path: 文件路径
            title: 标题（可选）
            slides: 幻灯片列表 [{"title": "标题", "body": "内容"}]
        """
        try:
            path = self._safe_path(params["path"], allow_create_parents=True)
        except PathSafetyError as e:
            return {"error": str(e)}

        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError:
            return {"error": "python-pptx 未安装，请运行: pip install python-pptx"}

        try:
            prs = Presentation()

            slide_data = params.get("slides", [])
            for slide_info in slide_data:
                # 使用空白布局，手动添加文本框，避免布局兼容性问题
                blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
                slide = prs.slides.add_slide(blank_layout)

                # 手动添加标题文本框
                title_text = slide_info.get("title", "")
                if title_text:
                    from pptx.util import Inches, Pt
                    title_box = slide.shapes.add_textbox(
                        Inches(0.5), Inches(0.3), Inches(9), Inches(1)
                    )
                    tf = title_box.text_frame
                    tf.text = title_text
                    for para in tf.paragraphs:
                        for run in para.runs:
                            run.font.size = Pt(28)
                            run.font.bold = True

                # 手动添加内容文本框
                body_text = slide_info.get("body", "")
                if body_text:
                    from pptx.util import Inches, Pt
                    body_box = slide.shapes.add_textbox(
                        Inches(0.5), Inches(1.5), Inches(9), Inches(5)
                    )
                    body_box.text_frame.text = body_text

            path.parent.mkdir(parents=True, exist_ok=True)
            prs.save(str(path))

            logger.info(f"创建 PPT: {path} ({len(slide_data)} 张幻灯片)")
            return {
                "path": str(path),
                "created": True,
                "slide_count": len(slide_data),
            }
        except Exception as e:
            logger.error(f"创建 PPT 失败: {e}")
            return {"error": f"创建失败: {e}"}

    # ------------------------------------------------------------------
    # PDF 文本提取
    # ------------------------------------------------------------------

    async def _read_pdf_text(self, params: dict) -> dict:
        """提取 PDF 文件中的文本内容。

        Params:
            path: 文件路径
            max_pages: 最大页数（默认 100）
        """
        try:
            path = self._safe_path(params["path"])
        except PathSafetyError as e:
            return {"error": str(e)}

        if not path.exists():
            return {"error": f"文件不存在: {path}"}

        try:
            import subprocess

            max_pages = params.get("max_pages", 100)

            # 优先使用 PyMuPDF (fitz)，回退到 pdftotext
            try:
                import fitz  # PyMuPDF

                doc = fitz.open(str(path))
                pages: list[dict[str, Any]] = []
                for i, page in enumerate(doc):
                    if i >= max_pages:
                        break
                    text = page.get_text()
                    pages.append({"page": i + 1, "text": text, "length": len(text)})
                doc.close()

                total_text = "\n".join(p["text"] for p in pages)
                logger.info(f"读取 PDF (PyMuPDF): {path} ({len(pages)} 页)")
                return {
                    "path": str(path),
                    "pages": pages,
                    "page_count": len(pages),
                    "total_text_length": len(total_text),
                }
            except ImportError:
                pass

            # 回退方案：使用 pdfminer
            try:
                from pdfminer.high_level import extract_text

                text = extract_text(str(path), maxpages=max_pages)
                logger.info(f"读取 PDF (pdfminer): {path}")
                return {
                    "path": str(path),
                    "text": text,
                    "total_text_length": len(text),
                }
            except ImportError:
                return {
                    "error": "PDF 处理库未安装，请安装: pip install PyMuPDF 或 pip install pdfminer.six"
                }
        except Exception as e:
            logger.error(f"读取 PDF 失败: {e}")
            return {"error": f"读取失败: {e}"}
