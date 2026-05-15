"""生成 JavasAgent 项目介绍 PPT。

用 python-pptx 自动生成一份关于 JavasAgent 的介绍演示文稿。
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os

# 颜色主题
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)       # 深蓝背景
ACCENT = RGBColor(0x00, 0xD4, 0xFF)         # 科技蓝
ACCENT2 = RGBColor(0x7B, 0x61, 0xFF)        # 紫色
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xAA, 0xAA, 0xAA)
GREEN = RGBColor(0x00, 0xE6, 0x76)
ORANGE = RGBColor(0xFF, 0x9F, 0x43)


def set_slide_bg(slide, color):
    """设置幻灯片背景色。"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="微软雅黑"):
    """添加文本框。"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=WHITE, icon="▸"):
    """添加带图标的列表。"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{icon} {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "微软雅黑"
        p.space_after = Pt(8)
    return txBox


def add_title_bar(slide, title, subtitle=""):
    """添加顶部标题栏。"""
    # 装饰线
    line = slide.shapes.add_shape(
        1,  # 矩形
        Inches(0.5), Inches(0.4), Inches(0.08), Inches(0.6)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    add_text_box(slide, 0.8, 0.3, 8, 0.5, title, font_size=28, bold=True, color=ACCENT)
    if subtitle:
        add_text_box(slide, 0.8, 0.8, 8, 0.4, subtitle, font_size=14, color=GRAY)


def create_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ========== 封面 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白
    set_slide_bg(slide, BG_DARK)

    # 大标题
    add_text_box(slide, 1, 1.5, 11, 1.2, "JavasAgent",
                 font_size=60, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)

    # 副标题
    add_text_box(slide, 1, 3.0, 11, 0.8, "像贾维斯一样的 AI 智能体，接管你的电脑",
                 font_size=28, color=WHITE, alignment=PP_ALIGN.CENTER)

    # 描述
    add_text_box(slide, 2, 4.2, 9, 0.6,
                 "感知 → 规划 → 决策 → 执行 → 反馈  |  语音唤醒 · 视觉驱动 · 技能扩展",
                 font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

    # 底部信息
    add_text_box(slide, 1, 6.2, 11, 0.5,
                 "GitHub: github.com/liusuren123/JavasAgent  |  v0.1.0  |  2026.05",
                 font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)

    # ========== 目录 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_title_bar(slide, "目录", "CONTENTS")

    toc_items = [
        ("01", "项目愿景", "打造什么样的 AI 智能体"),
        ("02", "系统架构", "五层架构 + 感知闭环"),
        ("03", "核心能力", "语音 · 视觉 · 操作 · 技能"),
        ("04", "语音助手", "唤醒词 → VAD → STT → Agent → TTS"),
        ("05", "技能系统", "YAML 零代码扩展能力"),
        ("06", "后台服务", "常驻 · 托盘 · 热键 · 自启"),
        ("07", "技术栈", "Python + Ollama + Porcupine"),
        ("08", "项目状态", "已完成模块与测试覆盖"),
    ]

    for i, (num, title, desc) in enumerate(toc_items):
        y = 1.5 + i * 0.7
        add_text_box(slide, 1.0, y, 0.6, 0.4, num, font_size=22, bold=True, color=ACCENT)
        add_text_box(slide, 1.8, y, 3, 0.4, title, font_size=20, bold=True, color=WHITE)
        add_text_box(slide, 5.0, y, 7, 0.4, desc, font_size=14, color=GRAY)

    # ========== 项目愿景 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_title_bar(slide, "项目愿景", "VISION")

    add_text_box(slide, 1, 1.5, 11, 0.6,
                 '打造一个类似钢铁侠中"贾维斯"的 AI 智能体系统',
                 font_size=22, color=WHITE, alignment=PP_ALIGN.CENTER)

    visions = [
        ("💻", "软件开发", "写代码、调试、测试、部署"),
        ("📊", "办公自动化", "操作 Office 文档、处理邮件、日程管理"),
        ("🎨", "创意工具", "操作 Photoshop、Premiere 等 Adobe 软件"),
        ("🌐", "通用操作", "文件管理、浏览器操作、信息检索"),
    ]

    for i, (icon, title, desc) in enumerate(visions):
        x = 1.0 + i * 3.0
        # 卡片背景
        card = slide.shapes.add_shape(1, Inches(x), Inches(2.8), Inches(2.6), Inches(2.5))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0x25, 0x25, 0x40)
        card.line.fill.background()

        add_text_box(slide, x + 0.2, 3.0, 2.2, 0.5, icon, font_size=36,
                     color=WHITE, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + 0.2, 3.7, 2.2, 0.4, title, font_size=18,
                     bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + 0.2, 4.2, 2.2, 0.6, desc, font_size=13,
                     color=GRAY, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, 1, 5.8, 11, 0.5,
                 "核心原则：决策不明确 → 询问人类  |  指令明确 → 自主执行  |  渐进式能力扩展",
                 font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

    # ========== 系统架构 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_title_bar(slide, "系统架构", "ARCHITECTURE — 五层架构")

    layers = [
        ("用户交互层", "CLI · 语音助手 · 对话窗口 · 系统托盘", ACCENT),
        ("核心引擎层", "Planner 规划 · Executor 执行 · Decider 决策 · Scheduler 调度", RGBColor(0x4E, 0xC5, 0xF1)),
        ("感知层", "ScreenAnalyzer · OCR · VisionEye · TargetMatcher · TargetCache", ACCENT2),
        ("平台操作层", "WindowsPlatform · HumanHand · MotorController · AudioStream", GREEN),
        ("基础支撑层", "Config · Logger · LLM Client · Memory · SkillRegistry", ORANGE),
    ]

    for i, (name, desc, color) in enumerate(layers):
        y = 1.5 + i * 1.1
        # 层标签
        label = slide.shapes.add_shape(1, Inches(0.8), Inches(y), Inches(2.2), Inches(0.8))
        label.fill.solid()
        label.fill.fore_color.rgb = color
        label.line.fill.background()
        tf = label.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
        p.font.name = "微软雅黑"
        p.alignment = PP_ALIGN.CENTER

        # 描述
        add_text_box(slide, 3.3, y + 0.1, 9, 0.6, desc, font_size=15, color=GRAY)

    # ========== 核心能力 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_title_bar(slide, "核心能力", "CORE CAPABILITIES")

    capabilities = [
        ("🎤", "语音交互", [
            "唤醒词检测（Porcupine）",
            "VAD 语音活动检测（Silero）",
            "STT 语音识别（Whisper/Google）",
            "TTS 语音合成（Edge-TTS）",
            "打断机制 + 连续对话",
        ]),
        ("👁️", "视觉感知", [
            "屏幕截图分析",
            "OCR 文字识别",
            "GroundingDINO 零样本定位",
            "三级目标匹配 fallback",
            "视觉闭环控制",
        ]),
        ("🖱️", "桌面操作", [
            "拟人鼠标移动（贝塞尔曲线）",
            "拟人打字（有间隔有错字）",
            "键盘快捷键",
            "窗口管理",
            "文件操作",
        ]),
        ("🔧", "工具集", [
            "30+ 内置工具",
            "代码开发辅助",
            "浏览器控制",
            "邮件/日程管理",
            "YAML 技能零代码扩展",
        ]),
    ]

    for i, (icon, title, items) in enumerate(capabilities):
        x = 0.5 + i * 3.2
        # 卡片
        card = slide.shapes.add_shape(1, Inches(x), Inches(1.5), Inches(2.9), Inches(5.2))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0x25, 0x25, 0x40)
        card.line.fill.background()

        add_text_box(slide, x + 0.2, 1.7, 2.5, 0.5, icon, font_size=32,
                     color=WHITE, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + 0.2, 2.3, 2.5, 0.4, title, font_size=20,
                     bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)
        add_bullet_list(slide, x + 0.3, 3.0, 2.4, 3.5, items,
                        font_size=13, color=GRAY, icon="•")

    # ========== 语音助手 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_title_bar(slide, "语音助手", "VOICE ASSISTANT — 端到端语音管道")

    pipeline_steps = [
        ("WakeWord\n检测器", "Porcupine\n< 200ms", ACCENT),
        ("VAD\n语音检测", "Silero\n< 50ms", ACCENT2),
        ("STT\n语音识别", "Whisper\n0.5-2s", GREEN),
        ("Agent\n处理", "LLM\n1-3s", ORANGE),
        ("TTS\n语音合成", "Edge-TTS\n实时", RGBColor(0xE0, 0x5A, 0xFF)),
    ]

    for i, (name, detail, color) in enumerate(pipeline_steps):
        x = 1.0 + i * 2.3
        # 圆角矩形
        box = slide.shapes.add_shape(1, Inches(x), Inches(2.2), Inches(2.0), Inches(1.8))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0x25, 0x25, 0x40)
        box.line.color.rgb = color
        box.line.width = Pt(2)

        add_text_box(slide, x + 0.1, 2.4, 1.8, 0.6, name, font_size=16,
                     bold=True, color=color, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + 0.1, 3.2, 1.8, 0.6, detail, font_size=12,
                     color=GRAY, alignment=PP_ALIGN.CENTER)

        # 箭头
        if i < len(pipeline_steps) - 1:
            add_text_box(slide, x + 2.0, 2.8, 0.4, 0.4, "→", font_size=24,
                         color=ACCENT, alignment=PP_ALIGN.CENTER)

    # 特性列表
    features = [
        "三级唤醒词降级：Porcupine → OpenWakeWord → VAD 模拟",
        "打断机制：Agent 说话时用户开口 → 0.5s 内停止 TTS",
        "连续对话模式：唤醒后持续听，超时再回到唤醒等待",
        "免唤醒模式：直接对话，适合桌面专注场景",
        "CLI 命令：javas voice / javas voice --continuous / javas voice --no-wake",
    ]
    add_bullet_list(slide, 1, 4.5, 11, 2.5, features, font_size=14, color=GRAY, icon="✦")

    # ========== 技能系统 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_title_bar(slide, "技能系统", "SKILL SYSTEM — YAML 零代码扩展")

    # 左侧：YAML 示例
    yaml_code = """name: "Word 另存为 PDF"
triggers: ["word 转 pdf", "另存为 pdf"]
steps:
  - action: key_combo
    keys: "f12"
  - action: wait_text
    text: "另存为"
  - action: click_text
    text: "PDF"
  - action: key_combo
    keys: "enter\""""

    code_box = slide.shapes.add_shape(1, Inches(0.8), Inches(1.5), Inches(5.5), Inches(4.5))
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = RGBColor(0x0D, 0x0D, 0x1A)
    code_box.line.color.rgb = RGBColor(0x33, 0x33, 0x55)

    add_text_box(slide, 1.0, 1.6, 5.1, 0.3, "skills/office/word_save_pdf.yaml",
                 font_size=11, color=ACCENT)
    add_text_box(slide, 1.0, 2.0, 5.1, 3.8, yaml_code, font_size=13,
                 color=RGBColor(0xCC, 0xCC, 0xCC), font_name="Consolas")

    # 右侧：说明
    add_text_box(slide, 7, 1.5, 5.5, 0.4, "20 个原子操作",
                 font_size=22, bold=True, color=ACCENT)

    action_groups = [
        ("键盘", "key_combo · key_type · type_text"),
        ("鼠标", "click · double_click · drag · scroll · move_mouse"),
        ("视觉", "click_text(OCR) · click_icon(视觉定位)"),
        ("等待", "wait · wait_text"),
        ("断言", "assert_text · assert_screen"),
        ("控制流", "condition(条件) · loop(循环) · run_skill(嵌套)"),
    ]

    for i, (group, actions) in enumerate(action_groups):
        y = 2.2 + i * 0.6
        add_text_box(slide, 7, y, 2, 0.3, group, font_size=15, bold=True, color=GREEN)
        add_text_box(slide, 9, y, 4, 0.3, actions, font_size=12, color=GRAY)

    add_text_box(slide, 7, 5.8, 5.5, 0.4,
                 "✦ 无需写代码，只需描述操作步骤  ✦ 安全：不用 eval，循环上限 100",
                 font_size=12, color=GRAY)

    # ========== 后台服务 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_title_bar(slide, "后台常驻服务", "DAEMON SERVICE — 始终在线")

    daemon_features = [
        ("🖥️", "系统托盘", "pystray 托盘图标\n🟢监听 🟡处理 🔴异常\n右键菜单控制"),
        ("⌨️", "全局热键", "Ctrl+Alt+J 呼出对话\nCtrl+Alt+V 切换语音\nCtrl+Alt+S 停止任务"),
        ("📡", "IPC 通信", "Named Pipe\nCLI ↔ 后台服务\nJSON-RPC 2.0"),
        ("💬", "对话窗口", "tkinter 轻量 GUI\n输入框 + 输出区\n关闭=隐藏不退出"),
        ("🔄", "开机自启", "Windows 注册表\n用户登录自动启动\n无需手动打开"),
    ]

    for i, (icon, title, desc) in enumerate(daemon_features):
        x = 0.5 + i * 2.5
        card = slide.shapes.add_shape(1, Inches(x), Inches(1.5), Inches(2.3), Inches(3.0))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0x25, 0x25, 0x40)
        card.line.fill.background()

        add_text_box(slide, x + 0.1, 1.7, 2.1, 0.4, icon, font_size=28,
                     color=WHITE, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + 0.1, 2.2, 2.1, 0.4, title, font_size=17,
                     bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + 0.2, 2.8, 1.9, 1.5, desc, font_size=12,
                     color=GRAY, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, 1, 5.0, 11, 0.5,
                 "CLI 命令：javas service start | javas service stop | javas service install | javas service status",
                 font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)

    # ========== 技术栈 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_title_bar(slide, "技术栈", "TECH STACK")

    tech_categories = [
        ("语言 & 运行时", [
            "Python 3.11+",
            "asyncio 异步架构",
            "Ollama 本地 LLM",
        ]),
        ("LLM 层", [
            "Ollama qwen3.6（本地默认）",
            "智谱 GLM / OpenAI GPT（云端）",
            "多 provider 可切换",
        ]),
        ("语音层", [
            "Porcupine 唤醒词",
            "Silero VAD",
            "faster-whisper STT",
            "Edge-TTS 语音合成",
        ]),
        ("感知层", [
            "OpenCV 图像处理",
            "GroundingDINO 零样本定位",
            "PaddleOCR / Windows OCR",
            "ChromaDB 向量记忆",
        ]),
        ("桌面控制", [
            "pyautogui 键鼠控制",
            "pywin32 Win32 API",
            "pystray 系统托盘",
            "keyboard 全局热键",
        ]),
    ]

    for i, (cat, items) in enumerate(tech_categories):
        x = 0.5 + i * 2.5
        add_text_box(slide, x + 0.1, 1.5, 2.3, 0.4, cat, font_size=16,
                     bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)
        add_bullet_list(slide, x + 0.2, 2.1, 2.2, 3.5, items,
                        font_size=13, color=GRAY, icon="•")

    add_text_box(slide, 1, 5.5, 11, 0.4,
                 "测试：pytest 2022+ 测试用例  |  全绿  |  覆盖率 98.9%+",
                 font_size=14, color=GREEN, alignment=PP_ALIGN.CENTER)

    # ========== 项目状态 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_title_bar(slide, "项目状态", "PROJECT STATUS")

    # 左列
    left_items = [
        ("核心引擎", "✅ 完成", "Planner · Executor · Decider · Scheduler"),
        ("平台层", "✅ 完成", "截屏 · 鼠标 · 键盘 · 窗口管理"),
        ("拟人手部", "✅ 完成", "贝塞尔曲线 · 拟人打字 · 点击偏移"),
        ("感知层", "✅ 完成", "OCR · 目标匹配 · VisionEye · 闭环控制"),
        ("语音助手", "✅ 完成", "唤醒词 · VAD · STT · TTS · 管道"),
    ]

    for i, (name, status, desc) in enumerate(left_items):
        y = 1.5 + i * 0.9
        add_text_box(slide, 1, y, 2, 0.3, name, font_size=17, bold=True, color=WHITE)
        add_text_box(slide, 3.2, y, 1.2, 0.3, status, font_size=14, bold=True, color=GREEN)
        add_text_box(slide, 1, y + 0.35, 5, 0.3, desc, font_size=12, color=GRAY)

    # 右列
    right_items = [
        ("技能系统", "✅ 完成", "YAML 加载 · StepExecutor · 20 原语 · 12 预置技能"),
        ("后台服务", "✅ 完成", "系统托盘 · 热键 · IPC · 开机自启 · 对话窗口"),
        ("工具集", "✅ 完成", "30+ 工具模块 · 统一 execute() 接口"),
        ("记忆系统", "✅ 完成", "短期 · 长期(ChromaDB) · 知识库 · 用户偏好"),
        ("文档", "✅ 完成", "架构 · 技术栈 · 使用说明 · API 参考 · 贡献指南"),
    ]

    for i, (name, status, desc) in enumerate(right_items):
        y = 1.5 + i * 0.9
        add_text_box(slide, 7, y, 2, 0.3, name, font_size=17, bold=True, color=WHITE)
        add_text_box(slide, 9.2, y, 1.2, 0.3, status, font_size=14, bold=True, color=GREEN)
        add_text_box(slide, 7, y + 0.35, 5.5, 0.3, desc, font_size=12, color=GRAY)

    # 数字统计
    stats = [
        ("2022+", "测试用例"),
        ("100%", "测试通过"),
        ("30+", "工具模块"),
        ("12", "预置技能"),
        ("20", "原子操作"),
        ("50+", "Git Commits"),
    ]

    for i, (num, label) in enumerate(stats):
        x = 1.0 + i * 2.0
        add_text_box(slide, x, 6.0, 1.8, 0.4, num, font_size=28, bold=True,
                     color=ACCENT, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x, 6.5, 1.8, 0.3, label, font_size=12,
                     color=GRAY, alignment=PP_ALIGN.CENTER)

    # ========== 结尾 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, 1, 2.0, 11, 1.0, "JavasAgent",
                 font_size=48, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 3.2, 11, 0.6, "让 AI 真正操控你的电脑",
                 font_size=24, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 4.5, 11, 0.4,
                 "github.com/liusuren123/JavasAgent",
                 font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 5.5, 11, 0.4, "Thanks!",
                 font_size=32, bold=True, color=RGBColor(0x55, 0x55, 0x77),
                 alignment=PP_ALIGN.CENTER)

    # 保存
    output_dir = os.path.dirname(os.path.abspath(__file__))
    output_path = os.path.join(output_dir, "JavasAgent-Introduction.pptx")
    prs.save(output_path)
    print(f"PPT 已生成: {output_path}")
    return output_path


if __name__ == "__main__":
    create_ppt()
